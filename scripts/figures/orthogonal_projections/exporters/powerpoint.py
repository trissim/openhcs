"""PowerPoint backend for slide deck export."""

from pathlib import Path
from typing import Sequence

from pptx import Presentation
from pptx.util import Inches

from .layouts import SlideBox, TwoPanelLayoutComposer
from .models import DeckExportOptions, DeckExportResult, FigureAsset, SlideContentSpec


class PptxDeckExporter:
    """Writes slide definitions to a PowerPoint deck."""

    def __init__(self) -> None:
        self._layout_composer = TwoPanelLayoutComposer()

    def export(
        self,
        slides: Sequence[SlideContentSpec],
        options: DeckExportOptions,
    ) -> DeckExportResult:
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

        for slide_spec in slides:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            layout = self._layout_composer.compose(
                slide_spec,
                options.slide_layout,
                slide_width_inches=13.333,
                slide_height_inches=7.5,
            )
            self._add_textbox(slide, layout.title_box, slide_spec.title, 24, bold=True)
            if slide_spec.subtitle:
                self._add_textbox(slide, layout.subtitle_box, slide_spec.subtitle, 12)
            if slide_spec.notes:
                footer_text = " | ".join(slide_spec.notes)
                self._add_textbox(slide, layout.footer_box, footer_text, 8)
            for asset in slide_spec.assets[:2]:
                asset_box = layout.asset_boxes[asset.asset_type]
                self._add_picture(slide, asset.path, asset_box)

        options.output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(options.output_path))
        return DeckExportResult(
            output_path=options.output_path, slide_count=len(slides)
        )

    def _add_textbox(
        self, slide, box: SlideBox, text: str, font_size: int, bold: bool = False
    ) -> None:
        shape = slide.shapes.add_textbox(
            Inches(box.left),
            Inches(box.top),
            Inches(box.width),
            Inches(box.height),
        )
        paragraph = shape.text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = text
        run.font.size = Inches(font_size / 72.0)
        run.font.bold = bold

    def _add_picture(self, slide, path: Path, box: SlideBox) -> None:
        slide.shapes.add_picture(
            str(path),
            Inches(box.left),
            Inches(box.top),
            width=Inches(box.width),
            height=Inches(box.height),
        )
