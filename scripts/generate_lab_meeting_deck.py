"""Generate the CellProfiler/OpenHCS lab-meeting PowerPoint deck."""

from __future__ import annotations

from collections.abc import Sequence
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image


OUTPUT_PATH = Path("docs/openhcs_cellprofiler_lab_meeting.pptx")
FIGURE_DIR = Path("/tmp/openhcs_labmeeting_prelim/cppipe_figures_20260505_v7_compact")
PROJECTED_FIGURE_DIR = Path(
    "/tmp/openhcs_labmeeting_prelim/cppipe_figures_20260505_v7_projected_8sample"
)
SUMMARY_FIGURE_DIR = Path(
    "/tmp/openhcs_labmeeting_prelim/cppipe_figures_20260505_v7_summary_distribution"
)
FIGURE_SLIDE_IMAGE_WIDTH = Inches(12.25)
FIGURE_SLIDE_IMAGE_LEFT = Inches(0.52)
FIGURE_SLIDE_IMAGE_TOP = Inches(1.05)
EXAMPLE_IMAGE_DIR = Path("local/muecs_cp/images")
CP_OUTPUT_IMAGE_DIR = Path(
    "/tmp/openhcs_labmeeting_prelim/native_cp3_official_refreshed_current"
)
LOGO_IMAGE_DIR = Path("docs/pics")
CP_EXAMPLE_IMAGE_DIR = LOGO_IMAGE_DIR / "examples"

NAVY = RGBColor(23, 43, 77)
TEAL = RGBColor(15, 139, 141)
ORANGE = RGBColor(217, 95, 2)
SLATE = RGBColor(73, 80, 87)
LIGHT_BG = RGBColor(247, 249, 250)
PALE_TEAL = RGBColor(226, 244, 243)
PALE_ORANGE = RGBColor(255, 239, 224)
WHITE = RGBColor(255, 255, 255)
PURPLE = RGBColor(98, 76, 169)
BLUE = RGBColor(54, 101, 167)
GREEN = RGBColor(64, 145, 108)


@dataclass(frozen=True)
class BulletSlide:
    title: str
    subtitle: str | None = None
    bullets: tuple[str, ...] = ()
    footer: str | None = None


@dataclass(frozen=True)
class FigureSlide:
    title: str
    figure: Path
    caption: str


@dataclass(frozen=True)
class ImageTile:
    title: str
    path: Path


BULLET_SLIDES: tuple[BulletSlide, ...] = (
    BulletSlide(
        title="CellProfiler Is A Big Deal",
        subtitle="This is the tool a lot of image-analysis people already trust",
        bullets=(
            "It turns microscopy images into tables of measurements.",
            "It is free, open-source, and widely used in biology labs.",
            "NIGMS says CellProfiler has been cited in more than 15,000 scientific papers.",
            "So matching CellProfiler is a serious test, not a toy benchmark.",
        ),
        footer="Sources: NIGMS Biomedical Beat; Genome Biology article metrics.",
    ),
    BulletSlide(
        title="Why This Matters",
        subtitle="Faster is only useful if the biology stays the same",
        bullets=(
            "If the numbers change, the biological conclusion can change.",
            "So the first question is: do the OpenHCS results match CellProfiler?",
            "Only after that do we ask how much faster it runs.",
        ),
    ),
    BulletSlide(
        title="What OpenHCS Adds",
        subtitle="OpenHCS makes the trusted pipeline easier to run and connect",
        bullets=(
            "Load a CellProfiler pipeline.",
            "Run the same analysis inside OpenHCS.",
            "Send images and results to napari, Fiji/ImageJ, OMERO-style storage, and normal tables.",
            "Keep the same results, but make the workflow easier to scale.",
        ),
        footer="Source: OpenHCS README and benchmark runtime implementation.",
    ),
    BulletSlide(
        title="How We Tested It",
        subtitle="Keep the comparison simple and conservative",
        bullets=(
            "Run the pipeline in CellProfiler.",
            "Run the converted version in OpenHCS.",
            "Compare the output tables and images.",
            "Compare how long each run took on one core.",
        ),
    ),
    BulletSlide(
        title="What We Can Say Today",
        subtitle="This is the polished figure set for lab meeting",
        bullets=(
            "18 official CellProfiler example pipelines are plotted here.",
            "All 18 pass the result check.",
            "All 18 are at least 4x faster on one core.",
            "More pipelines are being added, but these are the clean figures for now.",
        ),
    ),
    BulletSlide(
        title="Main Takeaway",
        subtitle="Same trusted result shape, much faster runs",
        bullets=(
            "CellProfiler is a major scientific imaging tool, so matching it matters.",
            "OpenHCS keeps the CellProfiler-style outputs while running faster.",
            "The v7 figures show at least 4x speedup across all 18 plotted pipelines.",
            "OpenHCS also gives one place to connect CellProfiler, napari, Fiji/ImageJ, and OMERO-style data.",
        ),
    ),
)

FIGURE_SLIDES: tuple[FigureSlide, ...] = (
    FigureSlide(
        "Do The Results Still Match?",
        FIGURE_DIR / "cppipe_accuracy_zoom.png",
        "Zoomed axis confirms the result is not hiding a broad spread below perfect parity.",
    ),
    FigureSlide(
        "How Long Does It Take?",
        FIGURE_DIR / "cppipe_raw_seconds.png",
        "Absolute runtime: how long the same pipeline takes in native CellProfiler versus OpenHCS.",
    ),
    FigureSlide(
        "Same Runtime Chart, Log Scale",
        FIGURE_DIR / "cppipe_raw_seconds_log.png",
        "Log scale keeps short and long pipelines readable on one chart.",
    ),
    FigureSlide(
        "Speedup On One Core",
        FIGURE_DIR / "cppipe_speedup.png",
        "Every plotted pipeline clears the 4x speedup target on single-core, single-sample execution.",
    ),
    FigureSlide(
        "Same Speedup Chart, Log Scale",
        FIGURE_DIR / "cppipe_speedup_log.png",
        "Some pipelines are far beyond the minimum 4x threshold.",
    ),
    FigureSlide(
        "Speedup By Assay Type",
        SUMMARY_FIGURE_DIR / "cppipe_v7_assay_category_speedup_dots.png",
        "Assay-level summary of speedup across the plotted benchmark set.",
    ),
    FigureSlide(
        "Speedup By Assay Type, Log Scale",
        SUMMARY_FIGURE_DIR / "cppipe_v7_assay_category_speedup_dots_log.png",
        "Log-scale assay summary keeps the very large wins readable.",
    ),
    FigureSlide(
        "Speedup By Analysis Type",
        SUMMARY_FIGURE_DIR / "cppipe_v7_module_category_speedup_dots.png",
        "Module-level summary of speedup across the plotted benchmark set.",
    ),
    FigureSlide(
        "Speedup By Analysis Type, Log Scale",
        SUMMARY_FIGURE_DIR / "cppipe_v7_module_category_speedup_dots_log.png",
        "Log-scale analysis summary keeps the very large wins readable.",
    ),
    FigureSlide(
        "Measured Throughput: One Sample vs More Samples",
        SUMMARY_FIGURE_DIR
        / "cppipe_throughput_measured_selected_measured_throughput_speedup.png",
        "Measured 1-core single-sample runs versus measured 2-core 16-sample partial throughput.",
    ),
    FigureSlide(
        "Measured 3-Core Scaling Across 12 Wells",
        SUMMARY_FIGURE_DIR / "cppipe_well12_workers3_well_throughput_summary.png",
        "Measured 12-well, 3-core throughput summary across all 18 plotted pipelines.",
    ),
)

EXAMPLE_TILES: tuple[ImageTile, ...] = (
    ImageTile("Neural channels", EXAMPLE_IMAGE_DIR / "dapi-gfap_orig.png"),
    ImageTile("Processed DAPI", EXAMPLE_IMAGE_DIR / "dapi_processed.png"),
    ImageTile(
        "CellProfiler overlay",
        CP_OUTPUT_IMAGE_DIR
        / "ExampleHuman_ExampleHuman"
        / "ExampleHuman_ExampleHuman_native_cellprofiler"
        / "AS_09125_050116030001_D03f00d0_Overlay.png",
    ),
    ImageTile(
        "Yeast colonies",
        CP_OUTPUT_IMAGE_DIR
        / "ExampleYeastColonies_ExampleYeastColonies"
        / "ExampleYeastColonies_ExampleYeastColonies_native_cellprofiler"
        / "6-1_outlines.png",
    ),
)

CP_EXAMPLE_TILES: tuple[ImageTile, ...] = (
    ImageTile("Colocalization", CP_EXAMPLE_IMAGE_DIR / "examplecolocalization_orig_1.png"),
    ImageTile("Colocalization mask", CP_EXAMPLE_IMAGE_DIR / "examplecolocalization_masked_1.png"),
    ImageTile("Human nuclei", CP_EXAMPLE_IMAGE_DIR / "humannuclei_3.jpg"),
    ImageTile("Human color overlay", CP_EXAMPLE_IMAGE_DIR / "humancolornucleism_2.jpg"),
    ImageTile("Fly image", CP_EXAMPLE_IMAGE_DIR / "fruitflyimg_1.jpg"),
    ImageTile("Fly objects", CP_EXAMPLE_IMAGE_DIR / "fruitfly-coded_1.jpg"),
    ImageTile("Comet assay", CP_EXAMPLE_IMAGE_DIR / "examplecometimg_1.png"),
    ImageTile("Comet objects", CP_EXAMPLE_IMAGE_DIR / "examplecometid_4.png"),
    ImageTile("Tumor image", CP_EXAMPLE_IMAGE_DIR / "exampletumorimg_5.png"),
    ImageTile("Tumor objects", CP_EXAMPLE_IMAGE_DIR / "exampletumorid_1.png"),
    ImageTile("Yeast patches", CP_EXAMPLE_IMAGE_DIR / "exampleyeastpatch_1.png"),
    ImageTile("Yeast objects", CP_EXAMPLE_IMAGE_DIR / "exampleyeastpatchid_1.png"),
)


def main() -> None:
    missing_assets = [
        *(slide.figure for slide in FIGURE_SLIDES if not slide.figure.exists()),
        *(tile.path for tile in EXAMPLE_TILES if not tile.path.exists()),
        *(tile.path for tile in CP_EXAMPLE_TILES if not tile.path.exists()),
        *(
            logo_path
            for logo_path in (
        LOGO_IMAGE_DIR / "omero_logo.png",
        LOGO_IMAGE_DIR / "fiji_logo.png",
        LOGO_IMAGE_DIR / "cp_logo_text.png",
        LOGO_IMAGE_DIR / "napari_logo-300x300-1.png",
        LOGO_IMAGE_DIR / "cellprofilershuffle2_1.png",
            )
            if not logo_path.exists()
        ),
    ]
    if missing_assets:
        missing = "\n".join(str(path) for path in missing_assets)
        raise FileNotFoundError(f"Missing required asset(s):\n{missing}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_cellprofiler_context_slide(prs)
    add_integration_slide(prs)
    add_example_images_slide(prs)
    for slide in FIGURE_SLIDES:
        add_figure_slide(prs, slide)
    add_bullet_slide(prs, BULLET_SLIDES[4])
    add_bullet_slide(prs, BULLET_SLIDES[5])
    add_sources_slide(prs)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(OUTPUT_PATH)


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    return slide


def add_title(slide, text: str, *, top: float = 0.35, font_size: int = 30) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(12.2), Inches(0.6))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(font_size)
    p.font.color.rgb = NAVY


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.08), Inches(12.0), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(8.5)
    p.font.color.rgb = SLATE


def add_title_slide(prs: Presentation) -> None:
    slide = blank_slide(prs)
    add_accent_bar(slide)
    title = slide.shapes.add_textbox(Inches(0.75), Inches(1.35), Inches(8.89), Inches(0.73))
    frame = title.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = "Cell Profiler Support on OpenHCS"
    p.font.bold = True
    p.font.size = Pt(38)
    p.font.color.rgb = NAVY

    add_footer(slide, "Prepared for lab meeting | v7 figures | 25-pipeline and 33-pipeline figure sets in progress")


def add_bullet_slide(prs: Presentation, spec: BulletSlide) -> None:
    slide = blank_slide(prs)
    add_accent_bar(slide)
    add_title(slide, spec.title)
    if spec.subtitle:
        subtitle = slide.shapes.add_textbox(Inches(0.7), Inches(1.05), Inches(11.2), Inches(0.4))
        p = subtitle.text_frame.paragraphs[0]
        p.text = spec.subtitle
        p.font.size = Pt(17)
        p.font.color.rgb = TEAL

    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.45), Inches(4.8))
    frame = body.text_frame
    frame.clear()
    for index, bullet in enumerate(spec.bullets):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(21)
        p.font.color.rgb = NAVY
        p.space_after = Pt(13)
    if spec.footer:
        add_footer(slide, spec.footer)


def add_cellprofiler_context_slide(prs: Presentation) -> None:
    slide = blank_slide(prs)
    add_accent_bar(slide)
    add_title(slide, "CellProfiler History", font_size=30)

    subtitle = slide.shapes.add_textbox(Inches(0.7), Inches(1.05), Inches(5.65), Inches(0.45))
    p = subtitle.text_frame.paragraphs[0]
    p.text = "A trusted image-analysis workhorse since the early 2000s"
    p.font.size = Pt(17)
    p.font.color.rgb = TEAL

    bullets = (
        "Started in 2003; released as open-source software in 2005.",
        "Official site lists 61 public published-pipeline entries: 58 downloads plus 3 GitHub repos.",
        "Company example: Recursion publicly describes using CellProfiler plus CNNs in its image-processing pipeline.",
        "Industry links go back to the original paper: Merck and Novartis fellowships helped support early development.",
        "NIH RePORTER lists $6.16M for Broad/Carpenter image-based profiling grant R35GM122547 from FY2017-2026.",
    )
    body = slide.shapes.add_textbox(Inches(0.1575), Inches(1.856), Inches(6.82), Inches(5.47))
    frame = body.text_frame
    frame.clear()
    for index, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(19)
        p.font.color.rgb = NAVY
        p.space_after = Pt(10)

    slide.shapes.add_picture(
        str(LOGO_IMAGE_DIR / "cellprofilershuffle2_1.png"),
        Inches(6.82),
        Inches(1.08),
        width=Inches(6.1),
        height=Inches(6.1),
    )


def add_example_images_slide(prs: Presentation) -> None:
    slide = blank_slide(prs)
    add_accent_bar(slide)
    add_title(slide, "The Benchmarks Cover Real Imaging Jobs")
    left = 0.55
    top = 1.02
    tile_width = 3.05
    tile_height = 1.85
    x_gap = 0.12
    y_gap = 0.12
    for index, tile in enumerate(CP_EXAMPLE_TILES):
        row, column = divmod(index, 4)
        add_mosaic_tile(
            slide,
            tile,
            left=Inches(left + column * (tile_width + x_gap)),
            top=Inches(top + row * (tile_height + y_gap)),
            width=Inches(tile_width),
            height=Inches(tile_height),
        )


def add_integration_slide(prs: Presentation) -> None:
    slide = blank_slide(prs)
    add_accent_bar(slide)
    add_title(slide, "OpenHCS Sits Between The Tools We Already Use")
    slide.shapes.add_picture(
        str(LOGO_IMAGE_DIR / "omero_logo.png"),
        Inches(1.025),
        Inches(1.47),
        width=Inches(4.475),
        height=Inches(0.781),
    )
    slide.shapes.add_picture(
        str(LOGO_IMAGE_DIR / "fiji_logo.png"),
        Inches(7.25),
        Inches(1.25),
        width=Inches(2.75),
        height=Inches(2.75),
    )
    slide.shapes.add_picture(
        str(LOGO_IMAGE_DIR / "cp_logo_text.png"),
        Inches(0.75),
        Inches(2.565),
        width=Inches(5.775),
        height=Inches(1.436),
    )
    slide.shapes.add_picture(
        str(LOGO_IMAGE_DIR / "napari_logo-300x300-1.png"),
        Inches(9.75),
        Inches(0.95),
        width=Inches(3.75),
        height=Inches(3.75),
    )

    add_lane(
        slide,
        Inches(0.9),
        Inches(4.35),
        Inches(3.3),
        Inches(1.25),
        "Bring it in",
        "CellProfiler pipelines and microscope images.",
        PALE_ORANGE,
        ORANGE,
    )
    add_lane(
        slide,
        Inches(5.0),
        Inches(4.35),
        Inches(3.3),
        Inches(1.25),
        "Run it",
        "OpenHCS runs the work and keeps the outputs comparable.",
        PALE_TEAL,
        TEAL,
    )
    add_lane(
        slide,
        Inches(9.1),
        Inches(4.35),
        Inches(3.3),
        Inches(1.25),
        "Look at it",
        "Results can go to napari, OMERO-style storage, Fiji/ImageJ, and tables.",
        WHITE,
        NAVY,
    )
    note = slide.shapes.add_textbox(Inches(0.95), Inches(6.15), Inches(11.4), Inches(0.5))
    p = note.text_frame.paragraphs[0]
    p.text = "Message: OpenHCS is a bridge, not a replacement for the tools biologists already know."
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY


def add_architecture_slide(prs: Presentation) -> None:
    slide = blank_slide(prs)
    add_accent_bar(slide)
    add_title(slide, "What Is OpenHCS Doing?")

    y = 1.35
    x_positions = (0.55, 3.25, 6.0, 8.8, 11.0)
    labels = (
        ("CellProfiler\nfile", "The saved pipeline", PALE_ORANGE, ORANGE),
        ("Convert", "Turn it into OpenHCS steps", WHITE, TEAL),
        ("Run", "Process images and objects", PALE_TEAL, TEAL),
        ("Save", "Write images and tables", WHITE, NAVY),
        ("Check", "Compare against CellProfiler", PALE_ORANGE, ORANGE),
    )
    for x, (title, subtitle, fill, line) in zip(x_positions, labels, strict=True):
        add_node(slide, x, y, 1.9, 1.35, title, subtitle, fill, line)
    for x in (2.45, 5.15, 7.95, 10.35):
        add_arrow(slide, x, y + 0.53)

    add_lane(
        slide,
        Inches(0.72),
        Inches(4.15),
        Inches(3.25),
        Inches(1.25),
        "Tools people know",
        "CellProfiler + Fiji/ImageJ + napari",
        PALE_ORANGE,
        ORANGE,
    )
    add_lane(
        slide,
        Inches(4.35),
        Inches(4.15),
        Inches(3.95),
        Inches(1.25),
        "OpenHCS in the middle",
        "Same workflow shape, faster execution",
        PALE_TEAL,
        TEAL,
    )
    add_lane(
        slide,
        Inches(8.72),
        Inches(4.15),
        Inches(3.25),
        Inches(1.25),
        "The check",
        "CellProfiler output vs OpenHCS output",
        WHITE,
        NAVY,
    )

    note = slide.shapes.add_textbox(Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.65))
    p = note.text_frame.paragraphs[0]
    p.text = "Key point: OpenHCS is not just clicking Run in CellProfiler. It converts the pipeline and runs it directly."
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY


def add_figure_slide(prs: Presentation, spec: FigureSlide) -> None:
    slide = blank_slide(prs)
    add_title(slide, spec.title, top=0.25, font_size=27)
    slide.shapes.add_picture(
        str(spec.figure),
        FIGURE_SLIDE_IMAGE_LEFT,
        FIGURE_SLIDE_IMAGE_TOP,
        width=FIGURE_SLIDE_IMAGE_WIDTH,
    )


def add_sources_slide(prs: Presentation) -> None:
    add_bullet_slide(
        prs,
        BulletSlide(
            title="Sources And Claim Boundary",
            bullets=(
                "NIGMS Biomedical Beat: CellProfiler cited in more than 15,000 scientific papers.",
                "Genome Biology metrics: CellProfiler 2006 paper around 4,968 citations.",
                "CellProfiler site: project started in 2003 and lists 61 public published-pipeline entries.",
                "NIH RePORTER: Broad/Carpenter R35GM122547 award rows total $6.16M from FY2017-2026.",
                "Genome Biology 2006 paper: early support included Merck and Novartis fellowships.",
                "Google Cloud Recursion case study: Recursion uses CellProfiler plus CNNs in its imaging pipeline.",
                "Broad/CellProfiler pages: use cases, open-source framing, and sustained funding/support context.",
                "OpenHCS README: Fiji/ImageJ, napari, storage, and integration framing.",
                "Performance figures: v7 OpenHCS benchmark outputs from 18 official CellProfiler examples.",
            ),
            footer="Exact URLs are listed in docs/lab_meeting_cellprofiler_openhcs_slides.md",
        ),
    )


def add_node(slide, x: float, y: float, w: float, h: float, title: str, subtitle: str, fill, line) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = line
    p = frame.add_paragraph()
    p.text = subtitle
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9.5)
    p.font.color.rgb = SLATE


def add_lane(slide, left, top, width, height, title: str, body: str, fill, line) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = line
    p.alignment = PP_ALIGN.CENTER
    p = frame.add_paragraph()
    p.text = body
    p.font.size = Pt(11)
    p.font.color.rgb = SLATE
    p.alignment = PP_ALIGN.CENTER


def add_arrow(slide, x: float, y: float) -> None:
    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        Inches(x),
        Inches(y),
        Inches(0.5),
        Inches(0.28),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = TEAL
    arrow.line.color.rgb = TEAL


def add_logo_card(slide, name: str, subtitle: str, color, left, top) -> None:
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        Inches(1.75),
        Inches(1.45),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    mark = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        left + Inches(0.58),
        top + Inches(0.14),
        Inches(0.58),
        Inches(0.58),
    )
    mark.fill.solid()
    mark.fill.fore_color.rgb = color
    mark.line.color.rgb = color
    p = mark.text_frame.paragraphs[0]
    p.text = name[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE

    title_box = slide.shapes.add_textbox(
        left + Inches(0.08),
        top + Inches(0.78),
        Inches(1.59),
        Inches(0.26),
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = name
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = color

    subtitle_box = slide.shapes.add_textbox(
        left + Inches(0.08),
        top + Inches(1.08),
        Inches(1.59),
        Inches(0.25),
    )
    p = subtitle_box.text_frame.paragraphs[0]
    p.text = subtitle
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(7.5)
    p.font.color.rgb = SLATE


def add_image_tile(
    slide,
    tile: ImageTile,
    *,
    left,
    top,
    width,
    height,
    show_title: bool = False,
) -> None:
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = TEAL
    picture_top = top + Inches(0.08)
    picture_height = height - Inches(0.16 if not show_title else 0.46)
    add_picture_contained(
        slide,
        tile.path,
        left + Inches(0.08),
        picture_top,
        width - Inches(0.16),
        picture_height,
    )
    if show_title:
        box = slide.shapes.add_textbox(left + Inches(0.08), top + height - Inches(0.33), width - Inches(0.16), Inches(0.22))
        p = box.text_frame.paragraphs[0]
        p.text = tile.title
        p.font.bold = True
        p.font.size = Pt(9)
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER


def add_mosaic_tile(slide, tile: ImageTile, *, left, top, width, height) -> None:
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = WHITE
    picture = add_picture_cover(slide, tile.path, left, top, width, height)
    picture.line.color.rgb = WHITE

    label_height = Inches(0.28)
    label = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left,
        top + height - label_height,
        width,
        label_height,
    )
    label.fill.solid()
    label.fill.fore_color.rgb = NAVY
    label.fill.transparency = 18
    label.line.color.rgb = NAVY
    p = label.text_frame.paragraphs[0]
    p.text = tile.title
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(8.5)
    p.font.color.rgb = WHITE


def add_picture_contained(slide, path: Path, left, top, width, height) -> None:
    with Image.open(path) as image:
        image_width, image_height = image.size
    image_ratio = image_width / image_height
    box_ratio = width / height
    if image_ratio >= box_ratio:
        rendered_width = width
        rendered_height = width / image_ratio
        rendered_left = left
        rendered_top = top + (height - rendered_height) / 2
    else:
        rendered_height = height
        rendered_width = height * image_ratio
        rendered_left = left + (width - rendered_width) / 2
        rendered_top = top
    slide.shapes.add_picture(str(path), rendered_left, rendered_top, width=rendered_width, height=rendered_height)


def add_picture_cover(slide, path: Path, left, top, width, height):
    with Image.open(path) as image:
        image_width, image_height = image.size
    image_ratio = image_width / image_height
    box_ratio = width / height
    picture = slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    if image_ratio > box_ratio:
        visible_width = box_ratio / image_ratio
        crop = (1.0 - visible_width) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    elif image_ratio < box_ratio:
        visible_height = image_ratio / box_ratio
        crop = (1.0 - visible_height) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def add_accent_bar(slide) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(0.16),
        Inches(7.5),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.color.rgb = TEAL


if __name__ == "__main__":
    main()
